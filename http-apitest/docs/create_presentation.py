#!/usr/bin/env python3
"""Generate PowerPoint presentation for Control Plane Automation"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
ORANGE = RGBColor(255, 102, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(64, 64, 64)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullet_points, code_block=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area
    content_top = Inches(1.5)
    content_width = Inches(6) if code_block else Inches(12.333)

    content_box = slide.shapes.add_textbox(Inches(0.5), content_top, content_width, Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY
        p.space_after = Pt(12)

    # Code block if provided
    if code_block:
        code_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
        code_shape.fill.solid()
        code_shape.fill.fore_color.rgb = RGBColor(40, 44, 52)
        code_shape.line.fill.background()

        code_box = slide.shapes.add_textbox(Inches(7), Inches(1.7), Inches(5.6), Inches(5.1))
        tf = code_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code_block
        p.font.size = Pt(12)
        p.font.name = "Courier New"
        p.font.color.rgb = RGBColor(200, 200, 200)

    return slide

def add_architecture_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Docker Swarm Cluster Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Draw nodes
    node_width = Inches(3.5)
    node_height = Inches(3)
    y_pos = Inches(2.2)

    nodes = [
        ("Node 1 (Manager)", Inches(0.8)),
        ("Node 2 (Worker)", Inches(4.9)),
        ("Node 3 (Worker)", Inches(9))
    ]

    for node_name, x_pos in nodes:
        # Node box
        node_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, node_width, node_height)
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = LIGHT_BLUE
        node_shape.line.color.rgb = DARK_BLUE
        node_shape.line.width = Pt(2)

        # Node title
        node_title = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.2), node_width, Inches(0.5))
        tf = node_title.text_frame
        p = tf.paragraphs[0]
        p.text = node_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Control Plane container
        cp_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.3), y_pos + Inches(0.9), Inches(2.9), Inches(1.8))
        cp_shape.fill.solid()
        cp_shape.fill.fore_color.rgb = ORANGE
        cp_shape.line.fill.background()

        cp_text = slide.shapes.add_textbox(x_pos + Inches(0.3), y_pos + Inches(1.4), Inches(2.9), Inches(0.8))
        tf = cp_text.text_frame
        p = tf.paragraphs[0]
        p.text = "Control Plane\n:3000"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Arrow to nodes_config.json
    config_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(5.8), Inches(3.5), Inches(0.8))
    config_box.fill.solid()
    config_box.fill.fore_color.rgb = RGBColor(40, 44, 52)
    config_box.line.fill.background()

    config_text = slide.shapes.add_textbox(Inches(4.9), Inches(5.9), Inches(3.5), Inches(0.6))
    tf = config_text.text_frame
    p = tf.paragraphs[0]
    p.text = 'nodes_config.json\n["ip1", "ip2", "ip3"]'
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_pipeline_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Automation Pipeline Overview"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Pipeline stages
    stages = [
        ("Launch\nAWS VMs", "aws role", LIGHT_BLUE),
        ("Setup\nDocker Swarm", "docker_swarm", LIGHT_BLUE),
        ("Deploy\nControl Plane", "deploy_cp", ORANGE),
        ("Run\nTests", "http-apitest", RGBColor(0, 150, 0))
    ]

    stage_width = Inches(2.5)
    stage_height = Inches(2)
    y_pos = Inches(2.5)
    start_x = Inches(0.8)
    gap = Inches(0.5)

    for i, (title, subtitle, color) in enumerate(stages):
        x_pos = start_x + i * (stage_width + gap)

        # Stage box
        stage_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, stage_width, stage_height)
        stage_shape.fill.solid()
        stage_shape.fill.fore_color.rgb = color
        stage_shape.line.fill.background()

        # Stage title
        stage_title = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.4), stage_width, Inches(1))
        tf = stage_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Stage subtitle
        stage_sub = slide.shapes.add_textbox(x_pos, y_pos + Inches(1.4), stage_width, Inches(0.5))
        tf = stage_sub.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.alignment = PP_ALIGN.CENTER

        # Arrow (except last)
        if i < len(stages) - 1:
            arrow_x = x_pos + stage_width + Inches(0.1)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, y_pos + Inches(0.8), Inches(0.3), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()

    # Command box
    cmd_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5))
    cmd_shape.fill.solid()
    cmd_shape.fill.fore_color.rgb = RGBColor(40, 44, 52)
    cmd_shape.line.fill.background()

    cmd_text = slide.shapes.add_textbox(Inches(1), Inches(5.4), Inches(11.3), Inches(1.3))
    tf = cmd_text.text_frame
    p = tf.paragraphs[0]
    p.text = "# Run entire pipeline\ncd ansible/launch_aws_setup && ansible-playbook site.yaml\ncd control-plane && go test -v -tags=http_apitest ./http-apitest/tests/..."
    p.font.size = Pt(16)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(200, 200, 200)

    return slide

# Create slides
add_title_slide(prs,
    "Control Plane Automated Testing",
    "End-to-End Deployment & Testing Infrastructure")

add_pipeline_slide(prs)

add_content_slide(prs,
    "Infrastructure Provisioning (Ansible)",
    [
        "AWS EC2 instances with Rocky Linux",
        "Configurable node count via ec2_instance_count",
        "Automatic SSH key generation and distribution",
        "Docker CE installation on all nodes",
        "Docker Swarm cluster initialization",
        "Control Plane deployment from ghcr.io/pgedge/control-plane"
    ],
    code_block="""# site.yaml
- name: Launch AWS infrastructure
  roles:
    - aws

- name: Setup Docker Swarm
  roles:
    - docker_swarm

- name: Deploy Control Plane
  roles:
    - deploy_control_plane

- name: Save IPs for testing
  roles:
    - save_node_ips""")

add_architecture_slide(prs)

add_content_slide(prs,
    "HTTP API Test Framework",
    [
        "Go-based test suite with http_apitest build tag",
        "Flexible configuration: env var or JSON file",
        "Tests cluster initialization (GET /v1/cluster/init)",
        "Tests node joining (POST /v1/cluster/join)",
        "Supports 1-11 node clusters",
        "Auto-configured by Ansible via nodes_config.json"
    ],
    code_block="""// nodes_config.json
{
  "nodes": [
    "13.127.241.95",
    "13.201.116.205",
    "13.233.245.218"
  ]
}

// Run tests
go test -v -tags=http_apitest \\
  ./http-apitest/tests/...

// Or use environment variable
CP_NODE_IPS=ip1,ip2,ip3 \\
  go test -v -tags=http_apitest ...""")

add_content_slide(prs,
    "Summary & Benefits",
    [
        "Infrastructure as Code - Reproducible AWS environment",
        "Container Orchestration - Docker Swarm for multi-node deployment",
        "Automated Deployment - Control Plane from official releases",
        "Integrated Testing - Seamless handoff to test framework",
        "Single Command - Full pipeline with ansible-playbook site.yaml",
        "Scalable - Change ec2_instance_count to add more nodes"
    ])

# Save presentation
output_path = "/Users/usman/pge/repos/control-plane/http-apitest/docs/Control_Plane_Automation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
